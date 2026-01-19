"""PPTX Exporter - IR to PowerPoint"""
import logging
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.services.slide_studio.ir.schema import SlideIR, DeckIR
from app.services.slide_studio.export.strategy import ExportStrategy
from app.services.slide_studio.config import slide_studio_config

logger = logging.getLogger(__name__)


class PPTXExporter:
    """Exports DeckIR to PowerPoint"""
    
    def __init__(self):
        self.slide_width = slide_studio_config.SLIDE_WIDTH
        self.slide_height = slide_studio_config.SLIDE_HEIGHT
    
    def export_deck(self, deck_ir: DeckIR, output_path: str) -> str:
        """
        Export deck to PPTX file.
        
        Args:
            deck_ir: Deck IR
            output_path: Output file path
            
        Returns:
            Path to exported file
        """
        try:
            # Create presentation (16:9)
            prs = Presentation()
            prs.slide_width = Inches(self.slide_width / 96.0)  # Convert pixels to inches
            prs.slide_height = Inches(self.slide_height / 96.0)
            
            # Add slides
            for slide_ir in deck_ir.slides:
                self._add_slide(prs, slide_ir)
            
            # Save
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            prs.save(output_path)
            
            logger.info(f"Exported deck {deck_ir.deck_id} to {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"PPTX export error: {e}")
            raise
    
    def _add_slide(self, prs: Presentation, slide_ir: SlideIR):
        """Add a slide to presentation"""
        # Create slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background from theme
        self._add_background(slide, slide_ir)
        
        # Add slots
        for slot_id, slot in slide_ir.slots.items():
            if not slot.bbox:
                continue
            
            if slot.slot_type == "text":
                self._add_text_box(slide, slot, slot.bbox)
            elif slot.slot_type == "bullet":
                self._add_bullet_list(slide, slot, slot.bbox)
            elif slot.slot_type == "image":
                # Image slots will be handled as images (placeholder for now)
                pass
            elif slot.slot_type == "chart":
                # Chart slots will be handled as images (placeholder for now)
                pass
            elif slot.slot_type == "table":
                if ExportStrategy.should_export_as_image(slot):
                    # Complex table as image (placeholder)
                    pass
                else:
                    self._add_table(slide, slot, slot.bbox)
    
    def _add_text_box(self, slide, slot, bbox):
        """Add text box to slide"""
        from pptx.enum.text import MSO_ANCHOR
        
        left = Inches(bbox.x / 96.0)
        top = Inches(bbox.y / 96.0)
        width = Inches(bbox.w / 96.0)
        height = Inches(bbox.h / 96.0)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = slot.content
        
        # Apply style
        if slot.style:
            paragraph = text_frame.paragraphs[0]
            if slot.style.font_size:
                paragraph.font.size = Pt(slot.style.font_size)
            if slot.style.color:
                paragraph.font.color.rgb = self._hex_to_rgb(slot.style.color)
            if slot.style.align:
                align_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT
                }
                paragraph.alignment = align_map.get(slot.style.align, PP_ALIGN.LEFT)
    
    def _add_bullet_list(self, slide, slot, bbox):
        """Add bullet list to slide"""
        left = Inches(bbox.x / 96.0)
        top = Inches(bbox.y / 96.0)
        width = Inches(bbox.w / 96.0)
        height = Inches(bbox.h / 96.0)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add bullet items
        for i, item in enumerate(slot.items):
            if i == 0:
                p = text_frame.paragraphs[0]
                p.text = item
            else:
                p = text_frame.add_paragraph()
                p.text = item
            
            p.level = 0
            p.font.size = Pt(slot.style.font_size if slot.style and slot.style.font_size else 18)
    
    def _add_table(self, slide, slot, bbox):
        """Add table to slide"""
        from pptx.enum.shapes import MSO_SHAPE
        
        left = Inches(bbox.x / 96.0)
        top = Inches(bbox.y / 96.0)
        width = Inches(bbox.w / 96.0)
        height = Inches(bbox.h / 96.0)
        
        rows = len(slot.rows) + 1  # +1 for header
        cols = len(slot.headers)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set headers
        for i, header in enumerate(slot.headers):
            cell = table.cell(0, i)
            cell.text = header
        
        # Set rows
        for row_idx, row_data in enumerate(slot.rows, 1):
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_data)
    
    def _add_background(self, slide, slide_ir: SlideIR):
        """Add background image to slide"""
        try:
            from app.services.slide_studio.store.repo import slide_studio_repo
            from app.services.slide_studio.theme.registry import theme_registry
            from app.services.slide_studio.background.renderer import background_renderer
            import tempfile
            import os
            
            # Get deck IR to access theme
            deck_ir = None
            for deck_id in slide_studio_repo._deck_irs.keys():
                candidate = slide_studio_repo.get_deck_ir(deck_id)
                if candidate and any(s.slide_id == slide_ir.slide_id for s in candidate.slides):
                    deck_ir = candidate
                    break
            
            if deck_ir and deck_ir.theme_id:
                theme = theme_registry.get_theme(deck_ir.theme_id)
                variant_id = deck_ir.theme_variant_map.get(slide_ir.slide_type.value)
                
                if variant_id:
                    variant = theme_registry.get_variant_for_slide_type(
                        deck_ir.theme_id,
                        slide_ir.slide_type.value,
                        variant_id
                    )
                    
                    # Render CSS background to PNG
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        png_path = tmp.name
                    
                    background_renderer.render_css_to_png(variant.css, png_path)
                    
                    # Set as slide background
                    from pptx.util import Inches
                    slide.background.fill.solid()
                    # For image background, we need to use fill.picture() instead
                    # But python-pptx has limitations, so we'll add as a shape behind
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(self.slide_width / 96.0)
                    height = Inches(self.slide_height / 96.0)
                    
                    pic = slide.shapes.add_picture(png_path, left, top, width, height)
                    # Move to back
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(2, pic._element)
                    
                    # Cleanup
                    try:
                        os.unlink(png_path)
                    except:
                        pass
                        
        except Exception as e:
            logger.warning(f"Background addition failed: {e}, using default")
    
    def _hex_to_rgb(self, hex_color: str):
        """Convert hex color to RGB"""
        from pptx.dml.color import RGBColor
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )


# Singleton instance
pptx_exporter = PPTXExporter()
